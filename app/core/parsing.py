"""Document parsing for various file formats."""

import io
import logging
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

from bs4 import BeautifulSoup
import markdown

logger = logging.getLogger(__name__)


@dataclass
class ParsedPage:
    """Represents a page or section of a document."""

    page_number: int
    content: str
    metadata: dict = field(default_factory=dict)


@dataclass
class ParsedDocument:
    """Result of parsing a document."""

    pages: list[ParsedPage]
    total_pages: int
    metadata: dict = field(default_factory=dict)
    parser_used: Literal["native", "mistral_ocr"] = "native"

    @property
    def full_text(self) -> str:
        """Get full text content."""
        return "\n\n".join(page.content for page in self.pages)


class DocumentParser(ABC):
    """Abstract base class for document parsers."""

    @abstractmethod
    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """Parse document content."""
        pass

    @staticmethod
    def clean_text(text: str) -> str:
        """Clean and normalize text."""
        # Normalize whitespace
        text = re.sub(r"\s+", " ", text)
        # Remove control characters except newlines
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
        return text.strip()


class PDFParser(DocumentParser):
    """Parser for PDF files using pypdf."""

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        pages = []

        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            text = self.clean_text(text)
            if text:
                pages.append(ParsedPage(page_number=i, content=text))

        metadata = {}
        if reader.metadata:
            metadata = {
                "title": reader.metadata.get("/Title", ""),
                "author": reader.metadata.get("/Author", ""),
                "subject": reader.metadata.get("/Subject", ""),
            }

        return ParsedDocument(
            pages=pages,
            total_pages=len(reader.pages),
            metadata=metadata,
        )


class DOCXParser(DocumentParser):
    """Parser for DOCX files using python-docx."""

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        from docx import Document

        doc = Document(io.BytesIO(content))
        
        # DOCX doesn't have real pages, so we group by headings or paragraphs
        paragraphs = []
        current_section = []
        current_title = None
        
        for para in doc.paragraphs:
            text = self.clean_text(para.text)
            if not text:
                continue
                
            # Check if it's a heading
            if para.style.name.startswith("Heading"):
                if current_section:
                    paragraphs.append((current_title, "\n".join(current_section)))
                current_title = text
                current_section = []
            else:
                current_section.append(text)
        
        if current_section:
            paragraphs.append((current_title, "\n".join(current_section)))

        # If no structure found, treat as single page
        if not paragraphs:
            full_text = "\n".join(self.clean_text(p.text) for p in doc.paragraphs if p.text)
            pages = [ParsedPage(page_number=1, content=full_text)] if full_text else []
        else:
            pages = [
                ParsedPage(
                    page_number=i + 1,
                    content=content,
                    metadata={"section_title": title} if title else {},
                )
                for i, (title, content) in enumerate(paragraphs)
            ]

        # Extract core properties
        metadata = {}
        if doc.core_properties:
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "subject": doc.core_properties.subject or "",
            }

        return ParsedDocument(
            pages=pages,
            total_pages=len(pages),
            metadata=metadata,
        )


class PPTXParser(DocumentParser):
    """Parser for PPTX files using python-pptx."""

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        pages = []

        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(self.clean_text(shape.text))
            
            if texts:
                pages.append(ParsedPage(
                    page_number=i,
                    content="\n".join(texts),
                    metadata={"slide_number": i},
                ))

        return ParsedDocument(
            pages=pages,
            total_pages=len(prs.slides),
            metadata={},
        )


class HTMLParser(DocumentParser):
    """Parser for HTML files using BeautifulSoup."""

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        # Try to detect encoding
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        soup = BeautifulSoup(text, "lxml")
        
        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()

        # Get text content
        body = soup.find("body") or soup
        text_content = self.clean_text(body.get_text(separator=" "))

        # Extract metadata
        metadata = {}
        title_tag = soup.find("title")
        if title_tag:
            metadata["title"] = title_tag.get_text()

        meta_author = soup.find("meta", attrs={"name": "author"})
        if meta_author:
            metadata["author"] = meta_author.get("content", "")

        pages = [ParsedPage(page_number=1, content=text_content)] if text_content else []

        return ParsedDocument(
            pages=pages,
            total_pages=1,
            metadata=metadata,
        )


class MarkdownParser(DocumentParser):
    """Parser for Markdown files."""

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        # Decode content
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        # Convert to HTML then extract text
        html = markdown.markdown(text)
        soup = BeautifulSoup(html, "lxml")
        text_content = self.clean_text(soup.get_text(separator=" "))

        # Try to extract title from first heading
        metadata = {}
        first_heading = soup.find(["h1", "h2"])
        if first_heading:
            metadata["title"] = first_heading.get_text()

        pages = [ParsedPage(page_number=1, content=text_content)] if text_content else []

        return ParsedDocument(
            pages=pages,
            total_pages=1,
            metadata=metadata,
        )


class TextParser(DocumentParser):
    """Parser for plain text files."""

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        # Try common encodings
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                text = content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text = content.decode("utf-8", errors="replace")

        text_content = self.clean_text(text)
        pages = [ParsedPage(page_number=1, content=text_content)] if text_content else []

        return ParsedDocument(
            pages=pages,
            total_pages=1,
            metadata={"filename": filename},
        )


# Content type to parser mapping
PARSERS: dict[str, type[DocumentParser]] = {
    "application/pdf": PDFParser,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DOCXParser,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": PPTXParser,
    "text/html": HTMLParser,
    "text/markdown": MarkdownParser,
    "text/plain": TextParser,
}

# Extension fallback mapping
EXTENSION_TO_CONTENT_TYPE: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".html": "text/html",
    ".htm": "text/html",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".txt": "text/plain",
}


def get_parser(content_type: str, filename: str) -> DocumentParser:
    """Get appropriate parser for content type."""
    # Try content type first
    parser_class = PARSERS.get(content_type)
    
    if not parser_class:
        # Fallback to extension
        ext = Path(filename).suffix.lower()
        inferred_type = EXTENSION_TO_CONTENT_TYPE.get(ext)
        if inferred_type:
            parser_class = PARSERS.get(inferred_type)
    
    if not parser_class:
        # Default to text parser
        parser_class = TextParser
    
    return parser_class()


def parse_document(content: bytes, filename: str, content_type: str) -> ParsedDocument:
    """Parse document content using appropriate parser (sync, no OCR)."""
    parser = get_parser(content_type, filename)
    return parser.parse(content, filename)


async def parse_document_with_ocr(
    content: bytes,
    filename: str,
    content_type: str,
) -> ParsedDocument:
    """Parse document with optional Mistral OCR fallback for PDFs.

    Heuristic: run native extraction first. If the extracted text is shorter
    than ``ocr_heuristic_min_text_chars``, re-parse with Mistral OCR.
    """
    from app.config import get_settings
    from app.services.document_ai.mistral_ocr import OCRProviderError, ocr_pdf

    settings = get_settings()
    is_pdf = content_type == "application/pdf" or Path(filename).suffix.lower() == ".pdf"

    # Non-PDF or OCR disabled → use native parser
    if not is_pdf or not settings.use_mistral_ocr or not settings.ocr_only_for_pdf:
        return parse_document(content, filename, content_type)

    # 1) Try native extraction first
    native_result = parse_document(content, filename, content_type)
    native_text_len = len(native_result.full_text.strip())

    if native_text_len >= settings.ocr_heuristic_min_text_chars:
        logger.info(
            "Native PDF extraction sufficient (%d chars) for %s",
            native_text_len,
            filename,
        )
        return native_result

    # 2) Native text too short → try OCR
    logger.info(
        "Native extraction only %d chars (threshold %d), attempting Mistral OCR for %s",
        native_text_len,
        settings.ocr_heuristic_min_text_chars,
        filename,
    )

    try:
        ocr_pages = await ocr_pdf(content, filename)
    except OCRProviderError:
        logger.warning("Mistral OCR failed for %s, falling back to native result", filename)
        return native_result

    if not ocr_pages:
        return native_result

    pages = [
        ParsedPage(
            page_number=p["page"],
            content=p["text"],
            metadata=p.get("meta") or {},
        )
        for p in ocr_pages
        if p["text"].strip()
    ]

    return ParsedDocument(
        pages=pages,
        total_pages=len(ocr_pages),
        metadata=native_result.metadata,
        parser_used="mistral_ocr",
    )
